"""
MIT License

Copyright (c) 2026 Mustafo Mahmadov (mustafomahmadov0061)

Permission is hereby granted, free of charge, to any person obtaining a copy
of this software and associated documentation files (the "Software"), to deal
in the Software without restriction, including without limitation the rights
to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
copies of the Software, and to permit persons to whom the Software is
furnished to do so, subject to the following conditions:

The above copyright notice and this permission notice shall be included in all
copies or substantial portions of the Software.

THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
SOFTWARE.
"""

import os
import sys
import uuid
import zipfile
import shutil
from pathlib import Path

# Save original user environment variables before any other imports/code might mutate them
_ORIGINAL_USERPROFILE = os.environ.get("USERPROFILE")
_ORIGINAL_HOME = os.environ.get("HOME")


import fitz  # PyMuPDF

import cv2
import numpy as np

TEMP_DIR = Path("temp")
TEMP_DIR.mkdir(exist_ok=True)
PREVIEW_DIR = TEMP_DIR / "previews"
PREVIEW_DIR.mkdir(exist_ok=True)
IMAGE_EXTENSIONS = ('.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.heic', '.webp')

APP_CACHE_DIR = Path(__file__).parent.parent / "app_cache"
APP_CACHE_DIR.mkdir(exist_ok=True)

# Standard A4 page dimensions in points (72 dpi)
A4_WIDTH = 595.0
A4_HEIGHT = 842.0

def run_office_com_convert(app_name, input_file, output_file, file_format):
    """
    Runs Office COM conversion, trying direct in-thread conversion first (required for frozen/PyInstaller environments)
    and falling back to an isolated subprocess if direct conversion fails (and not frozen).
    """
    import os
    import sys
    import subprocess
    
    # Temporarily restore original USERPROFILE/HOME variables for win32com/Office to work correctly
    curr_userprofile = os.environ.get("USERPROFILE")
    curr_home = os.environ.get("HOME")
    if _ORIGINAL_USERPROFILE:
        os.environ["USERPROFILE"] = _ORIGINAL_USERPROFILE
    if _ORIGINAL_HOME:
        os.environ["HOME"] = _ORIGINAL_HOME

    # Try direct in-process COM conversion first (which works in PyInstaller frozen environments)
    try:
        import win32com.client
        import pythoncom
        import winreg
        
        input_path = os.path.abspath(input_file)
        output_path = os.path.abspath(output_file)
        file_format = int(file_format)
        
        pythoncom.CoInitialize()
        try:
            if app_name == "Word":
                for version in ["15.0", "16.0"]:
                    try:
                        key = winreg.CreateKey(winreg.HKEY_CURRENT_USER, rf"SOFTWARE\Microsoft\Office\{version}\Word\Options")
                        winreg.SetValueEx(key, "DisableConvertPDFWarning", 0, winreg.REG_DWORD, 1)
                    except:
                        pass
                app = win32com.client.Dispatch("Word.Application")
                app.Visible = False
                app.DisplayAlerts = False
                doc = app.Documents.Open(input_path)
                doc.SaveAs(output_path, FileFormat=file_format)
                doc.Close(False)
                app.Quit()
            elif app_name == "Excel":
                app = win32com.client.Dispatch("Excel.Application")
                app.Visible = False
                app.DisplayAlerts = False
                wb = app.Workbooks.Open(input_path)
                if file_format == 57:
                    try:
                        for sheet in wb.Sheets:
                            try:
                                sheet.ResetAllPageBreaks()
                                sheet.PageSetup.PrintArea = False
                            except:
                                pass
                            try:
                                used_range = sheet.UsedRange
                                used_range.WrapText = True
                                for col_idx in range(1, used_range.Columns.Count + 1):
                                    if used_range.Columns(col_idx).ColumnWidth > 40:
                                        used_range.Columns(col_idx).ColumnWidth = 40
                            except Exception as e:
                                pass
                            sheet.PageSetup.Zoom = False
                            sheet.PageSetup.FitToPagesWide = 1
                            sheet.PageSetup.FitToPagesTall = 1
                            sheet.PageSetup.Orientation = 2
                    except:
                        pass
                    wb.ExportAsFixedFormat(0, output_path)
                else:
                    wb.SaveAs(output_path, FileFormat=file_format)
                wb.Close(False)
                app.Quit()
            elif app_name == "PowerPoint":
                app = win32com.client.Dispatch("PowerPoint.Application")
                pres = app.Presentations.Open(input_path, WithWindow=False)
                pres.SaveAs(output_path, file_format)
                pres.Close()
                app.Quit()
            
            print(f"[CONVERT] Direct Office COM conversion for {app_name} succeeded.")
            return True
        finally:
            pythoncom.CoUninitialize()
    except Exception as direct_err:
        print(f"[CONVERT] Direct Office COM conversion for {app_name} failed: {direct_err}")

    # Fallback to subprocess (only if not running inside PyInstaller frozen application)
    if getattr(sys, 'frozen', False):
        print(f"[CONVERT] Cannot use subprocess fallback in frozen application.")
        # Restore environment variables
        if curr_userprofile:
            os.environ["USERPROFILE"] = curr_userprofile
        elif "USERPROFILE" in os.environ:
            del os.environ["USERPROFILE"]
        if curr_home:
            os.environ["HOME"] = curr_home
        elif "HOME" in os.environ:
            del os.environ["HOME"]
        return False

    # Subprocess execution script
    script = f"""
import sys
import os
import win32com.client
import pythoncom

app_name = {repr(app_name)}
input_path = os.path.abspath(sys.argv[1])
output_path = os.path.abspath(sys.argv[2])
file_format = int(sys.argv[3])

try:
    pythoncom.CoInitialize()
    if app_name == "Word":
        import winreg
        for version in ["15.0", "16.0"]:
            try:
                key = winreg.CreateKey(winreg.HKEY_CURRENT_USER, rf"SOFTWARE\\Microsoft\\Office\\{{version}}\\Word\\Options")
                winreg.SetValueEx(key, "DisableConvertPDFWarning", 0, winreg.REG_DWORD, 1)
            except:
                pass
        app = win32com.client.Dispatch("Word.Application")
        app.Visible = False
        app.DisplayAlerts = False
        doc = app.Documents.Open(input_path)
        doc.SaveAs(output_path, FileFormat=file_format)
        doc.Close(False)
        app.Quit()
    elif app_name == "Excel":
        app = win32com.client.Dispatch("Excel.Application")
        app.Visible = False
        app.DisplayAlerts = False
        wb = app.Workbooks.Open(input_path)
        if file_format == 57:
            try:
                for sheet in wb.Sheets:
                    try:
                        sheet.ResetAllPageBreaks()
                        sheet.PageSetup.PrintArea = False
                    except:
                        pass
                    try:
                        used_range = sheet.UsedRange
                        used_range.WrapText = True
                        for col_idx in range(1, used_range.Columns.Count + 1):
                            if used_range.Columns(col_idx).ColumnWidth > 40:
                                used_range.Columns(col_idx).ColumnWidth = 40
                    except Exception as e:
                        pass
                    sheet.PageSetup.Zoom = False
                    sheet.PageSetup.FitToPagesWide = 1
                    sheet.PageSetup.FitToPagesTall = 1
                    sheet.PageSetup.Orientation = 2
            except:
                pass
            wb.ExportAsFixedFormat(0, output_path)
        else:
            wb.SaveAs(output_path, FileFormat=file_format)
        wb.Close(False)
        app.Quit()
    elif app_name == "PowerPoint":
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(input_path, WithWindow=False)
        pres.SaveAs(output_path, file_format)
        pres.Close()
        app.Quit()
    pythoncom.CoUninitialize()
    print("SUCCESS")
except Exception as e:
    print("ERROR:", e, file=sys.stderr)
    sys.exit(1)
"""

    try:
        res = subprocess.run(
            [sys.executable, "-c", script, str(input_file), str(output_file), str(file_format)],
            capture_output=True,
            text=True,
            check=True
        )
        return True
    except Exception as e:
        stderr = getattr(e, "stderr", "")
        print(f"[CONVERT] Isolated Office COM conversion for {app_name} failed: {e}. Stderr: {stderr}")
        return False
    finally:
        # Restore mutated values
        if curr_userprofile:
            os.environ["USERPROFILE"] = curr_userprofile
        elif "USERPROFILE" in os.environ:
            del os.environ["USERPROFILE"]
        if curr_home:
            os.environ["HOME"] = curr_home
        elif "HOME" in os.environ:
            del os.environ["HOME"]


def _normalize_page_to_a4(dest_doc: fitz.Document, src_doc: fitz.Document,
                          page_index: int, rotation: int = 0) -> None:
    """
    Copy a single page from *src_doc* into *dest_doc*, scaling it to fit A4.

    The source page content is rendered at high resolution and placed centred
    on a fresh A4 page, preserving aspect ratio.  This guarantees every page in
    the output PDF has identical dimensions (595 × 842 pt) regardless of the
    original page or image size.
    """
    src_page = src_doc.load_page(page_index)

    # Apply any user-requested rotation to the source before we measure it
    effective_rotation = (src_page.rotation + rotation) % 360
    src_page.set_rotation(effective_rotation)
    src_rect = src_page.rect  # after rotation

    # --- render source page to a high-res bitmap --------------------------
    # Use 2× zoom so text and thin lines stay crisp after downscale to A4.
    zoom = max(1.0, min(A4_WIDTH / src_rect.width, A4_HEIGHT / src_rect.height) * 2.0)
    mat = fitz.Matrix(zoom, zoom)
    pix = src_page.get_pixmap(matrix=mat, alpha=False)
    img_bytes = pix.tobytes("png")

    # --- create a blank A4 page -------------------------------------------
    new_page = dest_doc.new_page(width=A4_WIDTH, height=A4_HEIGHT)

    # --- compute centred fit rectangle ------------------------------------
    pw, ph = pix.width, pix.height
    scale = min(A4_WIDTH / pw, A4_HEIGHT / ph)
    dst_w = pw * scale
    dst_h = ph * scale
    x0 = (A4_WIDTH - dst_w) / 2.0
    y0 = (A4_HEIGHT - dst_h) / 2.0
    target_rect = fitz.Rect(x0, y0, x0 + dst_w, y0 + dst_h)

    new_page.insert_image(target_rect, stream=img_bytes)


def _resize_for_document_upscale(img: np.ndarray, scale: float = 2.0) -> np.ndarray:
    """Upscale document screenshots like online 2x image resizers, with a pixel cap."""
    h, w = img.shape[:2]
    max_pixels = 24_000_000
    target_pixels = h * w * scale * scale
    if target_pixels > max_pixels:
        scale = (max_pixels / float(h * w)) ** 0.5

    if scale <= 1.01:
        return img

    new_size = (max(1, int(round(w * scale))), max(1, int(round(h * scale))))
    return cv2.resize(img, new_size, interpolation=cv2.INTER_CUBIC)


def _enhance_document_for_ocr(img: np.ndarray) -> np.ndarray:
    """
    Make document screenshots readable without changing their structure.

    This intentionally behaves more like an online 2x image upscaler than a scanner
    cleanup filter: preserve colors, preserve thin spreadsheet grid lines, and sharpen
    text lightly. Heavy binarization makes OCR grids easier sometimes, but it looks bad
    and can create black blocks in blank Excel columns.
    """
    img = _resize_for_document_upscale(img, scale=2.0)

    lab = cv2.cvtColor(img, cv2.COLOR_BGR2LAB)
    l_channel, a_channel, b_channel = cv2.split(lab)
    blurred = cv2.GaussianBlur(l_channel, (0, 0), 0.75)
    l_channel = cv2.addWeighted(l_channel, 1.32, blurred, -0.32, 0)
    enhanced = cv2.cvtColor(cv2.merge([l_channel, a_channel, b_channel]), cv2.COLOR_LAB2BGR)

    # Remove only tiny JPG noise from pure white UI/page areas. Keep light gray Excel
    # row fills and grid lines intact.
    near_white = np.all(enhanced > 250, axis=2)
    enhanced[near_white] = (255, 255, 255)

    return enhanced


def enhance_image_quality(file_path: str) -> str:
    """
    Stand-alone image enhancement for user to manually apply before OCR.
    Uses a document-safe 2x visual upscale instead of photo/anime super-resolution.
    Returns the path to the enhanced image.
    """
    if not os.path.exists(file_path):
        return file_path

    img = cv2.imread(file_path)
    if img is None:
        return file_path

    # Generate a lossless PNG. Re-saving OCR inputs as JPG can add artifacts around text.
    file_id = str(uuid.uuid4())
    out_name = f"enhanced_{file_id}.png"
    out_path = TEMP_DIR / out_name

    print("[OCR] Enhancing image with document-safe 2x upscale.")
    img = _enhance_document_for_ocr(img)
    cv2.imwrite(str(out_path), img, [cv2.IMWRITE_PNG_COMPRESSION, 3])
    return str(out_path)



def ensure_pdf(file_path: str) -> str:
    """If the file is an image or office document, convert it to a temporary PDF and return the PDF path.
    Otherwise return the original file path."""
    import os
    file_lower = file_path.lower()
    if file_lower.endswith(IMAGE_EXTENSIONS):
        pdf_path = file_path + "_converted.pdf"
        if not os.path.exists(pdf_path):
            if file_lower.endswith('.heic'):
                try:
                    from PIL import Image
                    import pillow_heif
                    pillow_heif.register_heif_opener()
                    img = Image.open(file_path)
                    img.convert('RGB').save(file_path + '.jpg', 'JPEG')
                    img_doc = fitz.open(file_path + '.jpg')
                except Exception as e:
                    print(f"HEIC loading error: {e}")
                    img_doc = fitz.open(file_path)
            else:
                try:
                    img_doc = fitz.open(file_path)
                    pdf_bytes = img_doc.convert_to_pdf()
                    img_doc.close()
                    with open(pdf_path, "wb") as f:
                        f.write(pdf_bytes)
                    return pdf_path
                except Exception as e:
                    print(f"Image to PDF conversion with PyMuPDF failed: {e}")
                    from PIL import Image, ImageSequence

                    image = Image.open(file_path)
                    frames = []
                    for frame in ImageSequence.Iterator(image):
                        frames.append(frame.convert("RGB"))

                    if not frames:
                        raise RuntimeError("Görüntü dosyası PDF'e dönüştürülemedi.")

                    first_frame, extra_frames = frames[0], frames[1:]
                    first_frame.save(pdf_path, "PDF", resolution=100.0, save_all=True, append_images=extra_frames)
                    return pdf_path

            pdf_bytes = img_doc.convert_to_pdf()
            img_doc.close()
            with open(pdf_path, "wb") as f:
                f.write(pdf_bytes)
        return pdf_path
        
    elif file_lower.endswith(('.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx')):
        pdf_path = file_path + "_converted.pdf"
        if not os.path.exists(pdf_path):
            abs_in = os.path.abspath(file_path)
            abs_out = os.path.abspath(pdf_path)
            
            success = False
            if file_lower.endswith(('.doc', '.docx')):
                success = run_office_com_convert("Word", abs_in, abs_out, 17)
            elif file_lower.endswith(('.xls', '.xlsx')):
                success = run_office_com_convert("Excel", abs_in, abs_out, 57)
            elif file_lower.endswith(('.ppt', '.pptx')):
                success = run_office_com_convert("PowerPoint", abs_in, abs_out, 32)
                
            if not success:
                raise RuntimeError(f"Office to PDF conversion failed using Microsoft Office COM for: {file_path}")
        return pdf_path

    return file_path


def generate_thumbnails(file_path: str, original_filename: str) -> list:
    file_id = str(uuid.uuid4())
    thumbnails = []

    # Convert source files first so merge/split always receive a PDF path. This is
    # especially important for .tif/.tiff files, including multi-page TIFFs.
    actual_pdf_path = ensure_pdf(file_path)

    doc = fitz.open(actual_pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5))
        img_name = f"{file_id}_page_{page_num}.jpg"
        img_path = PREVIEW_DIR / img_name
        pix.pil_save(str(img_path), format="JPEG", optimize=True, quality=72)
        thumbnails.append({
            "id": f"{file_id}_{page_num}",
            "page_num": page_num,
            "url": f"/temp/previews/{img_name}",
            "original_file": actual_pdf_path,  # Use the converted PDF for merging later!
            "original_name": original_filename
        })
    doc.close()
    return thumbnails


def merge_pdfs(pages_data: list, output_filename: str = "birlestirilmis_belge.pdf") -> str:
    if output_filename == "birlestirilmis_belge.pdf":
        output_filename = f"birlestirilmis_belge_{uuid.uuid4().hex[:8]}.pdf"
    output_path = TEMP_DIR / output_filename
    merged_doc = fitz.open()

    for item in pages_data:
        pdf_file = ensure_pdf(item["original_file"])
        src_doc = fitz.open(pdf_file)
        rotation = item.get("rotation", 0)
        _normalize_page_to_a4(merged_doc, src_doc, item["page_num"], rotation)
        src_doc.close()

    merged_doc.save(str(output_path), garbage=4, deflate=True)
    merged_doc.close()
    return f"/temp/{output_filename}"


def images_to_single_pdf(file_paths: list[str], output_filename: str = "gorseller_birlestirilmis.pdf") -> str:
    if output_filename == "gorseller_birlestirilmis.pdf":
        output_filename = f"gorseller_birlestirilmis_{uuid.uuid4().hex[:8]}.pdf"

    output_path = TEMP_DIR / output_filename
    merged_doc = fitz.open()

    try:
        for file_path in file_paths:
            pdf_file = ensure_pdf(file_path)
            src_doc = fitz.open(pdf_file)
            try:
                for page_idx in range(len(src_doc)):
                    _normalize_page_to_a4(merged_doc, src_doc, page_idx)
            finally:
                src_doc.close()

        if len(merged_doc) == 0:
            raise RuntimeError("PDF'e eklenecek görüntü bulunamadı.")

        if output_path.exists():
            output_path.unlink()
        merged_doc.save(str(output_path), garbage=4, deflate=True, deflate_images=True)
    finally:
        merged_doc.close()

    return f"/temp/{output_filename}"


def split_pdf(pages_data: list, interval: int, output_zip: str = "bolunmus_belgeler.zip") -> dict:
    zip_path = TEMP_DIR / output_zip
    
    merged_doc = fitz.open()
    for item in pages_data:
        pdf_file = ensure_pdf(item["original_file"])
        src_doc = fitz.open(pdf_file)
        rotation = item.get("rotation", 0)
        _normalize_page_to_a4(merged_doc, src_doc, item["page_num"], rotation)
        src_doc.close()

    parts_info = []

    with zipfile.ZipFile(zip_path, 'w') as zipf:
        start = 0
        part = 1
        while start < len(merged_doc):
            end = min(start + interval - 1, len(merged_doc) - 1)
            new_doc = fitz.open()
            new_doc.insert_pdf(merged_doc, from_page=start, to_page=end)

            part_name = f"parca_{part}.pdf"
            part_path = TEMP_DIR / part_name
            new_doc.save(str(part_path), garbage=4, deflate=True)
            new_doc.close()

            zipf.write(part_path, part_name)
            # DO NOT remove so the UI can download them individually
            # os.remove(part_path)

            size_kb = part_path.stat().st_size / 1024
            parts_info.append({
                "filename": part_name,
                "url": f"/temp/{part_name}",
                "size_kb": round(size_kb, 2)
            })

            start = end + 1
            part += 1

    merged_doc.close()
    return {
        "zip_url": f"/temp/{output_zip}",
        "parts": parts_info
    }


def _find_ghostscript() -> str | None:
    """Locate the bundled Ghostscript console executable.

    Search order:
      1. _internal/gs/bin/gswin64c.exe  (bundled – PyInstaller or dev tree)
      2. gswin64c / gswin32c on system PATH
    Returns the absolute path or None.
    """
    # Bundled location (works for both dev and frozen builds)
    if getattr(sys, "frozen", False):
        base = Path(sys._MEIPASS)
    else:
        base = Path(__file__).parent.parent

    for candidate in [
        base / "_internal" / "gs" / "bin" / "gswin64c.exe",
        base / "_internal" / "gs" / "bin" / "gswin32c.exe",
        base / "gs" / "bin" / "gswin64c.exe",
    ]:
        if candidate.exists():
            return str(candidate)

    # Fall back to system PATH
    import shutil as _shutil
    for name in ("gswin64c", "gswin32c", "gs"):
        found = _shutil.which(name)
        if found:
            return found
    return None


def _ghostscript_compress(gs_exe: str, input_path: str, output_path: str,
                          preset: str, extra_args: list[str] | None = None) -> bool:
    """Run Ghostscript to rewrite a PDF with the given PDFSETTINGS preset.

    Returns True on success, False on failure.
    """
    import subprocess
    cmd = [
        gs_exe,
        "-dNOPAUSE", "-dBATCH", "-dQUIET",
        "-sDEVICE=pdfwrite",
        "-dCompatibilityLevel=1.4",
        f"-dPDFSETTINGS={preset}",
    ]
    if extra_args:
        cmd.extend(extra_args)
    cmd += [f"-sOutputFile={output_path}", input_path]

    try:
        creationflags = 0
        if os.name == 'nt':
            import subprocess
            creationflags = getattr(subprocess, 'CREATE_NO_WINDOW', 0x08000000)
        
        result = subprocess.run(cmd, capture_output=True, timeout=300, creationflags=creationflags)
        return result.returncode == 0 and os.path.exists(output_path)
    except Exception as e:
        print(f"Ghostscript error: {e}")
        return False


def compress_pdf(file_path: str, level: str, output_filename: str = "sikistirilmis_belge.pdf") -> str:
    if output_filename == "sikistirilmis_belge.pdf":
        output_filename = f"sikistirilmis_belge_{uuid.uuid4().hex[:8]}.pdf"
    pdf_file = ensure_pdf(file_path)
    output_path = TEMP_DIR / output_filename
    original_size = os.path.getsize(pdf_file)
    original_size_kb = original_size / 1024

    candidates = []
    candidate_base = f"{Path(output_filename).stem}_{uuid.uuid4().hex[:6]}"

    gs_exe = _find_ghostscript()

    # ── Ghostscript compression ──────────────────────────────────────────
    #
    #  Common flags added to ALL GS runs for maximum structural compression:
    #    -dSubsetFonts=true       → embed only used glyphs (critical for Arabic/CJK)
    #    -dCompressFonts=true     → Flate-encode font data
    #    -dCompressPages=true     → Flate-encode page content streams
    #    -dDetectDuplicateImages  → deduplicate identical images
    #    -dOptimize=true          → linearize & compact object table
    #
    #  Strategy: run MULTIPLE GS configurations, keep the smallest output.
    #  Different PDFs respond to different presets:
    #    - text-heavy (fonts) → /ebook with font subsetting wins
    #    - image-heavy (scans) → /screen with downsampling wins
    #
    GS_COMMON = [
        "-dSubsetFonts=true",
        "-dCompressFonts=true",
        "-dCompressPages=true",
        "-dDetectDuplicateImages=true",
        "-dOptimize=true",
        "-dFastWebView=false",
        "-dAutoFilterColorImages=true",
        "-dAutoFilterGrayImages=true",
        "-dDownsampleColorImages=true",
        "-dDownsampleGrayImages=true",
        "-dDownsampleMonoImages=true",
        "-dColorImageDownsampleType=/Bicubic",
        "-dGrayImageDownsampleType=/Bicubic",
        "-dMonoImageDownsampleType=/Bicubic",
    ]

    if gs_exe:
        # Define multiple GS passes to try per level
        gs_passes = {
            "low": [
                {"preset": "/ebook", "extra": GS_COMMON + [
                    "-dColorImageResolution=150",
                    "-dGrayImageResolution=150",
                    "-dMonoImageResolution=300",
                ]},
            ],
            "medium": [
                {"preset": "/ebook", "extra": GS_COMMON + [
                    "-dColorImageResolution=120",
                    "-dGrayImageResolution=120",
                    "-dMonoImageResolution=200",
                ]},
                {"preset": "/screen", "extra": GS_COMMON + [
                    "-dColorImageResolution=100",
                    "-dGrayImageResolution=100",
                    "-dMonoImageResolution=200",
                ]},
            ],
            "high": [
                # Pass A: /ebook with aggressive image downsampling (best for text-heavy)
                {"preset": "/ebook", "extra": GS_COMMON + [
                    "-dColorImageResolution=72",
                    "-dGrayImageResolution=72",
                    "-dMonoImageResolution=150",
                ]},
                # Pass B: /screen (best for image-heavy / scanned PDFs)
                {"preset": "/screen", "extra": GS_COMMON + [
                    "-dColorImageResolution=72",
                    "-dGrayImageResolution=72",
                    "-dMonoImageResolution=150",
                ]},
                # Pass C: /ebook with moderate settings (quality fallback)
                {"preset": "/ebook", "extra": GS_COMMON + [
                    "-dColorImageResolution=96",
                    "-dGrayImageResolution=96",
                    "-dMonoImageResolution=200",
                ]},
            ],
        }

        passes = gs_passes.get(level, gs_passes["medium"])
        for idx, gs_cfg in enumerate(passes):
            gs_path = TEMP_DIR / f"{candidate_base}_gs{idx}.pdf"
            if _ghostscript_compress(
                gs_exe, str(pdf_file), str(gs_path),
                preset=gs_cfg["preset"], extra_args=gs_cfg["extra"],
            ):
                gs_size = os.path.getsize(gs_path)
                # Only keep if it actually reduced size
                if gs_size < original_size:
                    candidates.append(gs_path)
                    print(f"[COMPRESS] GS pass {idx} ({gs_cfg['preset']}): "
                          f"{gs_size/1024:.1f} KB ({(1 - gs_size/original_size)*100:.1f}% saved)")
                else:
                    print(f"[COMPRESS] GS pass {idx} ({gs_cfg['preset']}): "
                          f"{gs_size/1024:.1f} KB — LARGER than original, discarding")
                    try:
                        gs_path.unlink()
                    except Exception:
                        pass
    else:
        print("[COMPRESS] Ghostscript not found, using PyMuPDF fallback")

    # ── Pass A: PyMuPDF STRUCTURAL cleanup (best for text-heavy PDFs) ────
    #
    #  This pass focuses on:
    #    - garbage=4: merge duplicate objects + compact xref table
    #    - subset_fonts(): remove unused glyphs (CRITICAL for Arabic/CJK)
    #    - deflate everything: fonts, streams, pages
    #    - use_objstms: pack objects into compressed object streams
    #    - NO image rewriting (preserves quality, fast for text PDFs)
    #
    structural_path = TEMP_DIR / f"{candidate_base}_structural.pdf"
    try:
        doc = fitz.open(pdf_file)
        # Analyse document: detect text vs image ratio
        sample_pages = min(10, len(doc))
        total_img_area = 0
        total_page_area = 0
        for i in range(sample_pages):
            page = doc[i]
            rect = page.rect
            total_page_area += rect.width * rect.height
            for img in page.get_images(full=True):
                try:
                    xref = img[0]
                    bbox_list = page.get_image_rects(xref)
                    for bbox in bbox_list:
                        total_img_area += bbox.width * bbox.height
                except Exception:
                    pass
        img_ratio = total_img_area / total_page_area if total_page_area else 0
        is_text_heavy = img_ratio < 0.3
        print(f"[COMPRESS] Document analysis: {len(doc)} pages, "
              f"image coverage {img_ratio*100:.1f}%, "
              f"{'TEXT-HEAVY' if is_text_heavy else 'IMAGE-HEAVY'}")

        try:
            doc.subset_fonts()
            print("[COMPRESS] Font subsetting completed")
        except Exception as e:
            print(f"[COMPRESS] Font subsetting skipped: {e}")

        doc.save(
            str(structural_path), garbage=4,
            deflate=True, deflate_images=True, deflate_fonts=True,
            clean=True, preserve_metadata=False,
            use_objstms=1, compression_effort=100,
        )
        doc.close()
        struct_size = os.path.getsize(structural_path)
        if struct_size < original_size:
            candidates.append(structural_path)
            print(f"[COMPRESS] PyMuPDF structural: {struct_size/1024:.1f} KB "
                  f"({(1 - struct_size/original_size)*100:.1f}% saved)")
        else:
            print(f"[COMPRESS] PyMuPDF structural: {struct_size/1024:.1f} KB — LARGER, discarding")
            try:
                structural_path.unlink()
            except Exception:
                pass
    except Exception as e:
        print(f"[COMPRESS] PyMuPDF structural failed: {e}")

    # ── Pass B: PyMuPDF image rewrite (best for image-heavy PDFs) ────────
    rewrite_profiles = {
        "low":    {"dpi_threshold": 250, "dpi_target": 200, "quality": 80, "gray": False, "garbage": 2},
        "medium": {"dpi_threshold": 150, "dpi_target": 110, "quality": 50, "gray": False, "garbage": 4},
        "high":   {"dpi_threshold": 100, "dpi_target": 96,  "quality": 50, "gray": False, "garbage": 4},
    }
    rw = rewrite_profiles.get(level, rewrite_profiles["medium"])

    rewrite_path = TEMP_DIR / f"{candidate_base}_rewrite.pdf"
    try:
        doc = fitz.open(pdf_file)
        doc.rewrite_images(
            dpi_threshold=rw["dpi_threshold"], dpi_target=rw["dpi_target"],
            quality=rw["quality"], lossy=True, lossless=True, bitonal=True,
            color=True, gray=True, set_to_gray=rw["gray"],
        )
        try:
            doc.subset_fonts()
        except Exception:
            pass
        doc.save(
            str(rewrite_path), garbage=rw["garbage"],
            deflate=True, deflate_images=True, deflate_fonts=True,
            clean=True, preserve_metadata=False,
            use_objstms=1, compression_effort=100,
        )
        doc.close()
        rw_size = os.path.getsize(rewrite_path)
        if rw_size < original_size:
            candidates.append(rewrite_path)
            print(f"[COMPRESS] PyMuPDF rewrite: {rw_size/1024:.1f} KB "
                  f"({(1 - rw_size/original_size)*100:.1f}% saved)")
        else:
            print(f"[COMPRESS] PyMuPDF rewrite: {rw_size/1024:.1f} KB — LARGER, discarding")
            try:
                rewrite_path.unlink()
            except Exception:
                pass
    except Exception as e:
        print(f"[COMPRESS] PyMuPDF rewrite failed: {e}")

    # ── Pass C: Two-pass hybrid (GS result → PyMuPDF structural cleanup) ─
    #
    #  If GS produced any candidate, run PyMuPDF structural cleanup on the
    #  SMALLEST GS output. This often squeezes another 20-40% because:
    #    - GS optimizes images/streams but leaves bloated font tables
    #    - PyMuPDF subset_fonts + garbage=4 trims the rest
    #
    gs_candidates = [c for c in candidates if "_gs" in c.name]
    if gs_candidates:
        smallest_gs = min(gs_candidates, key=lambda p: os.path.getsize(p))
        hybrid_path = TEMP_DIR / f"{candidate_base}_hybrid.pdf"
        try:
            doc = fitz.open(smallest_gs)
            try:
                doc.subset_fonts()
            except Exception:
                pass
            doc.save(
                str(hybrid_path), garbage=4,
                deflate=True, deflate_images=True, deflate_fonts=True,
                clean=True, preserve_metadata=False,
                use_objstms=1, compression_effort=100,
            )
            doc.close()
            hybrid_size = os.path.getsize(hybrid_path)
            if hybrid_size < original_size:
                candidates.append(hybrid_path)
                print(f"[COMPRESS] Hybrid (GS+PyMuPDF): {hybrid_size/1024:.1f} KB "
                      f"({(1 - hybrid_size/original_size)*100:.1f}% saved)")
            else:
                print(f"[COMPRESS] Hybrid: {hybrid_size/1024:.1f} KB — LARGER, discarding")
                try:
                    hybrid_path.unlink()
                except Exception:
                    pass
        except Exception as e:
            print(f"[COMPRESS] Hybrid pass failed: {e}")

    # ── Pass D: Two-pass chained (structural → image rewrite) ────────────
    #
    #  For maximum compression: take the structural cleanup result
    #  and THEN do image rewriting on it. Double savings.
    #
    if structural_path.exists() and os.path.getsize(structural_path) < original_size:
        chained_path = TEMP_DIR / f"{candidate_base}_chained.pdf"
        try:
            doc = fitz.open(structural_path)
            doc.rewrite_images(
                dpi_threshold=rw["dpi_threshold"], dpi_target=rw["dpi_target"],
                quality=rw["quality"], lossy=True, lossless=True, bitonal=True,
                color=True, gray=True, set_to_gray=rw["gray"],
            )
            doc.save(
                str(chained_path), garbage=4,
                deflate=True, deflate_images=True, deflate_fonts=True,
                clean=True, preserve_metadata=False,
                use_objstms=1, compression_effort=100,
            )
            doc.close()
            chained_size = os.path.getsize(chained_path)
            if chained_size < original_size:
                candidates.append(chained_path)
                print(f"[COMPRESS] Chained (struct+rewrite): {chained_size/1024:.1f} KB "
                      f"({(1 - chained_size/original_size)*100:.1f}% saved)")
            else:
                try:
                    chained_path.unlink()
                except Exception:
                    pass
        except Exception as e:
            print(f"[COMPRESS] Chained pass failed: {e}")

    if not candidates:
        shutil.copy2(pdf_file, output_path)
        return f"/temp/{output_filename}"
    # ── Pick the smallest candidate ──────────────────────────────────────
    best_path = min(candidates, key=lambda p: os.path.getsize(p))
    best_size = os.path.getsize(best_path)
    best_name = best_path.name

    if output_path.exists():
        output_path.unlink()

    # SAFETY NET: never return a file larger than the original
    if best_size >= original_size:
        shutil.copy2(pdf_file, output_path)
    else:
        os.replace(best_path, output_path)

    # Clean up leftover candidates
    for c in candidates:
        try:
            if c.exists() and c != output_path:
                c.unlink()
        except Exception:
            pass

    compressed_size_kb = os.path.getsize(output_path) / 1024
    saved_pct = round(max(0, (1 - compressed_size_kb / original_size_kb) * 100), 1) if original_size_kb else 0

    # Determine engine used
    if "_gs" in best_name and "_hybrid" not in best_name:
        engine = "ghostscript"
    elif "_hybrid" in best_name:
        engine = "ghostscript+pymupdf"
    elif "_chained" in best_name:
        engine = "pymupdf-chained"
    else:
        engine = "pymupdf"

    print(f"[COMPRESS] ✓ BEST: {engine} → {compressed_size_kb:.1f} KB ({saved_pct}% saved)")

    return {
        "url": f"/temp/{output_filename}",
        "size_kb": round(compressed_size_kb, 1),
        "original_size_kb": round(original_size_kb, 1),
        "saved_percent": saved_pct,
        "preserved_structure": "_gs" not in best_name or level != "high",
        "max_compression_mode": level == "high",
        "compression_profile": level,
        "engine": engine,
    }


def process_conversion(file_path: str, task_id: str) -> str:
    """Handles various conversion tasks including OCR."""
    import zipfile
    import shutil
    filename = Path(file_path).stem

    if task_id in ["img_to_excel", "pdf_to_excel"]:
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_name = f"{filename}_ocr_{timestamp}.xlsx"
        output_path = TEMP_DIR / output_name
        print(f"[CONVERT] Task: {task_id}, File: {file_path}")
        print(f"[CONVERT] Output: {output_path}")
        try:
            # If the input is a PDF, render the first page to an image first
            actual_image_path = file_path
            if file_path.lower().endswith('.pdf'):
                print(f"[CONVERT] Input is PDF, rendering to high-res image...")
                doc = fitz.open(file_path)
                page = doc.load_page(0)
                # High resolution render for best OCR accuracy
                pix = page.get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
                actual_image_path = str(TEMP_DIR / f"{filename}_page0_hires.png")
                pix.save(actual_image_path)
                doc.close()
                print(f"[CONVERT] Rendered to: {actual_image_path}")

            print(f"[CONVERT] Calling extract_table_to_excel...")
            # Use our custom high-accuracy table extraction engine
            from backend.table_ocr import extract_table_to_excel
            model_dir = str((APP_CACHE_DIR / "paddleocr").absolute())
            num_rows, num_cols = extract_table_to_excel(
                actual_image_path, str(output_path), model_storage_dir=model_dir
            )
            print(f"[CONVERT] SUCCESS: {num_rows} tables found")

        except Exception as e:
            import traceback
            print(f"[CONVERT] ERROR: {e}")
            traceback.print_exc()
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            err_msg = "".join(ch for ch in str(e) if ch.isprintable() or ch in (' ', '\t'))
            ws["A1"] = f"OCR Hatasi: {err_msg[:200]}"
            wb.save(str(output_path))

        # Read back the saved Excel file for preview data
        preview_data = None
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(output_path), data_only=True)
            ws = wb.active
            preview_data = []
            for row in ws.iter_rows(values_only=True):
                preview_data.append([str(c) if c is not None else "" for c in row])
        except Exception as e:
            print(f"[CONVERT] Preview data error: {e}")

        return {"url": f"/temp/{output_name}", "preview": preview_data}

    elif task_id in ["pdf_to_word", "pdf_to_docx"]:
        pdf_file = ensure_pdf(file_path)
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_name = f"{filename}_{timestamp}.docx"
        output_path = TEMP_DIR / output_name
        
        success = False
        try:
            from pdf2docx import Converter
            cv = Converter(pdf_file)
            cv.convert(str(output_path), start=0, end=None)
            cv.close()
            if output_path.exists() and output_path.stat().st_size > 0:
                success = True
        except Exception as e:
            print(f"[CONVERT] pdf2docx conversion failed: {e}. Trying Microsoft Word fallback...")
            
        if not success:
            try:
                abs_in = os.path.abspath(pdf_file)
                abs_out = os.path.abspath(output_path)
                if run_office_com_convert("Word", abs_in, abs_out, 16):
                    if output_path.exists() and output_path.stat().st_size > 0:
                        success = True
            except Exception as ex:
                print(f"[CONVERT] Word PDF-to-Word fallback failed: {ex}")
                
        return f"/temp/{output_name}"

    elif task_id == "pdf_to_jpg":
        pdf_file = ensure_pdf(file_path)
        output_name = f"{filename}_jpg_sayfalar.zip"
        output_path = TEMP_DIR / output_name
        first_preview_url = None
        pages_urls = []

        try:
            import io
            from PIL import Image
            doc = fitz.open(pdf_file)
            with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zipf:
                for page_num, page in enumerate(doc, start=1):
                    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                    img_name = f"{filename}_sayfa_{page_num}.jpg"
                    img_path = TEMP_DIR / img_name
                    
                    # PyInstaller safe Pillow save
                    img_data = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_data))
                    img.convert("RGB").save(str(img_path), "JPEG", quality=92)
                    
                    zipf.write(img_path, img_name)
                    pages_urls.append(f"/temp/{img_name}")

                    if first_preview_url is None:
                        preview_name = f"{uuid.uuid4().hex}_preview.jpg"
                        preview_path = PREVIEW_DIR / preview_name
                        shutil.copy2(img_path, preview_path)
                        first_preview_url = f"/temp/previews/{preview_name}"
            doc.close()
        except Exception as e:
            print(f"PDF to JPG error: {e}")

        return {"url": f"/temp/{output_name}", "preview": first_preview_url, "pages": pages_urls}

    elif task_id == "pdf_to_ppt":
        pdf_file = ensure_pdf(file_path)
        output_name = f"{filename}.pptx"
        output_path = TEMP_DIR / output_name
        try:
            from pptx import Presentation
            from pptx.util import Inches
            import io
            from PIL import Image

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_slide_layout = prs.slide_layouts[6]

            doc = fitz.open(pdf_file)
            for page in doc:
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                img_path = TEMP_DIR / f"{uuid.uuid4().hex}.jpg"
                
                # PyInstaller safe Pillow save
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                img.convert("RGB").save(str(img_path), "JPEG", quality=92)

                slide = prs.slides.add_slide(blank_slide_layout)
                page_ratio = page.rect.width / page.rect.height
                slide_ratio = prs.slide_width / prs.slide_height
                if page_ratio >= slide_ratio:
                    width = prs.slide_width
                    height = int(width / page_ratio)
                    left = 0
                    top = int((prs.slide_height - height) / 2)
                else:
                    height = prs.slide_height
                    width = int(height * page_ratio)
                    left = int((prs.slide_width - width) / 2)
                    top = 0

                slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
                try:
                    img_path.unlink()
                except Exception:
                    pass
            doc.close()
            prs.save(str(output_path))
        except Exception as e:
            print(f"PDF to PPT error: {e}")

        return f"/temp/{output_name}"

    elif task_id == "pdf_to_odt":
        pdf_file = ensure_pdf(file_path)
        output_name = f"{filename}.odt"
        output_path = TEMP_DIR / output_name
        
        success = False
        import tempfile
        sys_temp = Path(tempfile.gettempdir())
        
        # Try direct Word conversion first as it maintains high-fidelity formatting
        try:
            temp_odt_path = sys_temp / f"temp_{uuid.uuid4().hex}.odt"
            abs_in = os.path.abspath(pdf_file)
            abs_out = os.path.abspath(temp_odt_path)
            if run_office_com_convert("Word", abs_in, abs_out, 23): # wdFormatOpenDocumentText = 23
                if temp_odt_path.exists():
                    import shutil
                    shutil.move(str(temp_odt_path), str(output_path))
                    if output_path.exists() and output_path.stat().st_size > 0:
                        success = True
        except Exception as e:
            print(f"[CONVERT] Direct Word PDF-to-ODT failed: {e}. Trying pdf2docx fallback...")
            
        if not success:
            try:
                docx_path = sys_temp / f"temp_{uuid.uuid4().hex}.docx"
                from pdf2docx import Converter
                cv = Converter(pdf_file)
                cv.convert(str(docx_path), start=0, end=None)
                cv.close()
                
                temp_odt_path = sys_temp / f"temp_{uuid.uuid4().hex}.odt"
                abs_docx = os.path.abspath(docx_path)
                abs_odt = os.path.abspath(temp_odt_path)
                if run_office_com_convert("Word", abs_docx, abs_odt, 23):
                    if temp_odt_path.exists():
                        import shutil
                        shutil.move(str(temp_odt_path), str(output_path))
                        if output_path.exists() and output_path.stat().st_size > 0:
                            success = True
                
                try:
                    docx_path.unlink()
                except:
                    pass
            except Exception as e:
                print(f"[CONVERT] pdf2docx fallback ODT conversion failed: {e}")
                
        return f"/temp/{output_name}"

    elif task_id == "pdf_to_ods":
        import tempfile
        sys_temp = Path(tempfile.gettempdir())
        xlsx_path = sys_temp / f"temp_{uuid.uuid4().hex}.xlsx"
        output_name = f"{filename}.ods"
        output_path = TEMP_DIR / output_name
        
        success = False
        try:
            from backend.table_ocr import extract_table_to_excel
            actual_image_path = file_path
            temp_hires_path = None
            if file_path.lower().endswith('.pdf'):
                doc = fitz.open(file_path)
                page = doc.load_page(0)
                pix = page.get_pixmap(matrix=fitz.Matrix(3.0, 3.0))
                temp_hires_path = TEMP_DIR / f"{filename}_page0_hires_ods.png"
                actual_image_path = str(temp_hires_path)
                pix.save(actual_image_path)
                doc.close()
                
            model_dir = str((APP_CACHE_DIR / "paddleocr").absolute())
            extract_table_to_excel(actual_image_path, str(xlsx_path), model_storage_dir=model_dir)
            
            if temp_hires_path and temp_hires_path.exists():
                try:
                    temp_hires_path.unlink()
                except:
                    pass
            
            if os.path.exists(xlsx_path):
                temp_ods_path = sys_temp / f"temp_{uuid.uuid4().hex}.ods"
                abs_in = os.path.abspath(xlsx_path)
                abs_out = os.path.abspath(temp_ods_path)
                if run_office_com_convert("Excel", abs_in, abs_out, 60): # 60 = xlOpenDocumentSpreadsheet
                    if temp_ods_path.exists():
                        import shutil
                        shutil.move(str(temp_ods_path), str(output_path))
                
                try:
                    os.remove(xlsx_path)
                except:
                    pass
                if output_path.exists() and output_path.stat().st_size > 0:
                    success = True
        except Exception as e:
            print(f"PDF to ODS error: {e}")
            
        if not success:
            try:
                import openpyxl
                wb = openpyxl.Workbook()
                ws = wb.active
                ws["A1"] = "Dönüştürme hatası veya tablolar bulunamadı."
                wb.save(str(xlsx_path))
                wb.close()
                
                temp_ods_path = sys_temp / f"temp_{uuid.uuid4().hex}.ods"
                abs_in = os.path.abspath(xlsx_path)
                abs_out = os.path.abspath(temp_ods_path)
                if run_office_com_convert("Excel", abs_in, abs_out, 60):
                    if temp_ods_path.exists():
                        import shutil
                        shutil.move(str(temp_ods_path), str(output_path))
                
                try:
                    os.remove(xlsx_path)
                except:
                    pass
            except Exception as ex:
                print(f"ODS fallback failed: {ex}")
                
        return f"/temp/{output_name}"

    elif task_id == "pdf_to_odp":
        pdf_file = ensure_pdf(file_path)
        output_name = f"{filename}.odp"
        output_path = TEMP_DIR / output_name
        
        import tempfile
        sys_temp = Path(tempfile.gettempdir())
        pptx_path = sys_temp / f"temp_{uuid.uuid4().hex}.pptx"
        try:
            from pptx import Presentation
            from pptx.util import Inches
            import io
            from PIL import Image
            
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank_slide_layout = prs.slide_layouts[6]
            
            doc = fitz.open(pdf_file)
            temp_images = []
            for page in doc:
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                img_path = sys_temp / f"{uuid.uuid4().hex}.png"
                
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                img.save(str(img_path), "PNG")
                temp_images.append(img_path)
                
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
                
            doc.close()
            prs.save(str(pptx_path))
            
            # Clean up temp images
            for img_path in temp_images:
                try:
                    img_path.unlink()
                except:
                    pass
            
            temp_odp_path = sys_temp / f"temp_{uuid.uuid4().hex}.odp"
            abs_in = os.path.abspath(pptx_path)
            abs_out = os.path.abspath(temp_odp_path)
            if run_office_com_convert("PowerPoint", abs_in, abs_out, 35): # ppSaveAsOpenDocumentPresentation = 35
                if temp_odp_path.exists():
                    import shutil
                    shutil.move(str(temp_odp_path), str(output_path))
            
            try:
                pptx_path.unlink()
            except:
                pass
        except Exception as e:
            print(f"PDF to ODP error: {e}")
            try:
                if pptx_path.exists():
                    pptx_path.unlink()
            except:
                pass
        return f"/temp/{output_name}"
        
    elif task_id == "pdf_to_txt":
        pdf_file = ensure_pdf(file_path)
        output_name = f"{filename}.txt"
        output_path = TEMP_DIR / output_name
        try:
            doc = fitz.open(pdf_file)
            with open(output_path, "w", encoding="utf-8") as f:
                for page in doc:
                    f.write(page.get_text() + "\n\n")
            doc.close()
        except Exception as e:
            print(f"PDF to TXT error: {e}")
        return f"/temp/{output_name}"

    elif task_id in ["heic_to_jpg", "jpg_to_heic", "webp_to_jpg", "jpg_to_webp", "tiff_to_jpg", "jpg_to_tiff"]:
        from PIL import Image
        import pillow_heif
        pillow_heif.register_heif_opener()
        
        out_ext = ".jpg"
        save_format = "JPEG"
        if task_id in ["heic_to_jpg", "webp_to_jpg", "tiff_to_jpg"]:
            out_ext = ".jpg"
            save_format = "JPEG"
        elif task_id == "jpg_to_heic":
            out_ext = ".heic"
            save_format = "HEIF"
        elif task_id == "jpg_to_webp":
            out_ext = ".webp"
            save_format = "WEBP"
        elif task_id == "jpg_to_tiff":
            out_ext = ".tiff"
            save_format = "TIFF"
            
        uid = uuid.uuid4().hex[:6]
        output_name = f"{filename}_donusturulen_{uid}{out_ext}"
        output_path = TEMP_DIR / output_name
        
        try:
            img = Image.open(file_path)
            if save_format == "JPEG" and img.mode in ("RGBA", "P", "LA"):
                img = img.convert("RGB")
            img.save(str(output_path), format=save_format)
            img.close()
        except Exception as e:
            print(f"Image conversion error ({task_id}): {e}")
            raise e
            
        return f"/temp/{output_name}"

    elif task_id == "enhance_image":
        enhanced_path = enhance_image_quality(file_path)
        import ntpath
        output_name = ntpath.basename(enhanced_path)
        return f"/temp/{output_name}"

    elif task_id == "make_iso":
        import zipfile
        import pycdlib
        
        uid = uuid.uuid4().hex[:6]
        output_name = f"{filename}_olusturulan_{uid}.iso"
        output_path = TEMP_DIR / output_name
        
        iso = pycdlib.PyCdlib()
        iso.new(interchange_level=3, joliet=True)
        
        temp_extract_dir = None
        try:
            is_zip = False
            try:
                if zipfile.is_zipfile(file_path):
                    is_zip = True
            except:
                pass
                
            if is_zip:
                temp_extract_dir = TEMP_DIR / f"make_iso_extract_{uid}"
                temp_extract_dir.mkdir(exist_ok=True)
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_extract_dir)
                
                joliet_to_iso9660_dir = { '/': '/' }
                dir_id_counter = 0
                file_counter = 0
                
                for root, dirs, files in os.walk(temp_extract_dir):
                    rel_dir = os.path.relpath(root, temp_extract_dir).replace('\\', '/')
                    if rel_dir == '.':
                        joliet_parent = '/'
                    else:
                        joliet_parent = '/' + rel_dir
                        
                    for d in dirs:
                        if joliet_parent == '/':
                            joliet_dir = '/' + d
                        else:
                            joliet_dir = joliet_parent + '/' + d
                            
                        dir_id_counter += 1
                        iso9660_parent_path = joliet_to_iso9660_dir[joliet_parent]
                        if iso9660_parent_path == '/':
                            iso9660_dir = f"/D{dir_id_counter}"
                        else:
                            iso9660_dir = f"{iso9660_parent_path}/D{dir_id_counter}"
                            
                        iso.add_directory(iso_path=iso9660_dir, joliet_path=joliet_dir)
                        joliet_to_iso9660_dir[joliet_dir] = iso9660_dir
                        
                    for f in files:
                        file_counter += 1
                        full_local_path = os.path.join(root, f)
                        
                        if joliet_parent == '/':
                            joliet_file = '/' + f
                        else:
                            joliet_file = joliet_parent + '/' + f
                            
                        iso9660_parent_path = joliet_to_iso9660_dir[joliet_parent]
                        if iso9660_parent_path == '/':
                            iso9660_file = f"/F{file_counter};1"
                        else:
                            iso9660_file = f"{iso9660_parent_path}/F{file_counter};1"
                            
                        iso.add_file(full_local_path, iso9660_file, joliet_path=joliet_file)
            else:
                import ntpath
                fname = ntpath.basename(file_path)
                joliet_file = '/' + fname
                iso9660_file = "/F1;1"
                iso.add_file(file_path, iso9660_file, joliet_path=joliet_file)
                
            iso.write(str(output_path))
            iso.close()
        except Exception as e:
            print(f"Make ISO error: {e}")
            raise e
        finally:
            if temp_extract_dir and temp_extract_dir.exists():
                import shutil
                try:
                    shutil.rmtree(temp_extract_dir)
                except:
                    pass
                    
        return f"/temp/{output_name}"

    elif task_id == "extract_iso":
        import zipfile
        import pycdlib
        
        uid = uuid.uuid4().hex[:6]
        output_name = f"{filename}_ayiklangan_{uid}.zip"
        output_path = TEMP_DIR / output_name
        
        extracted_paths = []
        try:
            iso = pycdlib.PyCdlib()
            iso.open(file_path)
            
            use_joliet = iso.has_joliet()
            facade = iso.get_joliet_facade() if use_joliet else iso.get_iso9660_facade()
            
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                walker = facade.walk(joliet_path='/') if use_joliet else facade.walk(iso_path='/')
                for root, dirs, files in walker:
                    for f in files:
                        full_path = root + '/' + f if root != '/' else '/' + f
                        out_temp = TEMP_DIR / f"extracted_{uuid.uuid4().hex}_{f}"
                        
                        with open(out_temp, 'wb') as out_f:
                            if use_joliet:
                                facade.get_file_from_iso_fp(out_f, joliet_path=full_path)
                            else:
                                facade.get_file_from_iso_fp(out_f, iso_path=full_path)
                        
                        clean_name = f.split(';')[0] if ';' in f else f
                        parent_dir = root.lstrip('/')
                        arcname = parent_dir + '/' + clean_name if parent_dir else clean_name
                        
                        zipf.write(out_temp, arcname)
                        extracted_paths.append(out_temp)
            iso.close()
        except Exception as e:
            print(f"Extract ISO error: {e}")
            raise e
        finally:
            for p in extracted_paths:
                try:
                    p.unlink()
                except:
                    pass
        return f"/temp/{output_name}"

    # Generic fallback — convert to PDF
    pdf_file = ensure_pdf(file_path)
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_name = f"{filename}_converted_{timestamp}.pdf"
    output_path = TEMP_DIR / output_name
    doc = fitz.open(pdf_file)
    doc.save(str(output_path))
    doc.close()
    return f"/temp/{output_name}"
